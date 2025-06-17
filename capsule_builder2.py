import os
import json
from datetime import datetime
from pathlib import Path
from pptx import Presentation
import docx
import requests

ACTIVITY_LOG_FILE = "activity_log.jsonl"
CAPSULE_DIR = "capsules"
GROQ_API_KEY = "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"
GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"
GROQ_MODEL = "meta-llama/llama-4-scout-17b-16e-instruct"

Path(CAPSULE_DIR).mkdir(exist_ok=True)

def load_activity_log():
    activities = []
    if os.path.exists(ACTIVITY_LOG_FILE):
        with open(ACTIVITY_LOG_FILE, "r", encoding="utf-8") as f:
            for line in f:
                try:
                    entry = json.loads(line.strip())
                    if entry.get("app") and entry.get("title"):
                        activities.append(entry)
                except json.JSONDecodeError:
                    continue
    return activities

def extract_code_files():
    code_files = []
    for file in os.listdir("."):
        if file.endswith(".py"):
            with open(file, "r", encoding="utf-8") as f:
                code = f.read()
                code_files.append(f"--- {file} (PY) ---\n{code}")
    return "\n\n".join(code_files)

def extract_word_content(filename):
    try:
        doc = docx.Document(filename)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    except Exception as e:
        return f"[Error reading {filename}]: {str(e)}"

def extract_ppt_content(filename):
    try:
        prs = Presentation(filename)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text.strip())
        return "\n".join(text_runs)
    except Exception as e:
        return f"[Error reading {filename}]: {str(e)}"

def summarize_context(context_text):
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {GROQ_API_KEY}"
    }

    system_prompt = (
        "You are a summarization assistant. Read the session context below and summarize the user's goals and actions in 3 sentences."
    )

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": context_text}
    ]

    response = requests.post(GROQ_API_URL, headers=headers, json={"model": GROQ_MODEL, "messages": messages})
    if response.status_code == 200:
        return response.json()["choices"][0]["message"]["content"]
    else:
        return f"[Groq API error {response.status_code}]: {response.text}"

def save_capsule(context, summary):
    timestamp = datetime.now().strftime("%Y-%m-%dT%H-%M")
    capsule_path = Path(CAPSULE_DIR) / f"capsule_{timestamp}.json"
    with open(capsule_path, "w", encoding="utf-8") as f:
        json.dump({"context": context, "summary": summary}, f, indent=2)
    print(f"\n✅ Capsule saved: {capsule_path}")

if __name__ == "__main__":
    print("🧠 Building new capsule with Word + PPT + Logs + Code...\n")

    # Load activities
    activities = load_activity_log()
    activity_summary = ""
    file_contexts = []

    for entry in activities:
        app = entry["app"]
        title = entry["title"]
        activity_summary += f"- [{entry['timestamp']}] Opened `{title}` using `{app}`\n"

        # Parse Word file if title matches
        if "Word" in title and "SPEECH" in title:
            file_contexts.append(f"--- SPEECH.docx ---\n{extract_word_content('SPEECH.docx')}")
        elif "PowerPoint" in title and "MCP_A2A" in title:
            file_contexts.append(f"--- MCP_A2A.pptx ---\n{extract_ppt_content('MCP_A2A.pptx')}")

    # Add code files
    code_block = extract_code_files()
    full_context = f"Here is an activity summary of the user's session:\n{activity_summary}\n\n" + "\n\n".join(file_contexts) + "\n\n" + code_block

    # Summarize + Save
    print("📝 Summary:")
    summary = summarize_context(full_context)
    print("", summary)
    save_capsule(full_context, summary)
